import os
from io import BytesIO
import pandas as pd
from pptx import Presentation, chart
from pptx.chart.axis import _BaseAxis
from pptx.chart.data import ChartData
from pptx.chart.plot import PlotFactory, PlotTypeInspector
from pptx.enum.chart import XL_CHART_TYPE, XL_MARKER_STYLE
from pptx.util import Inches
import openpyxl
from sqlalchemy import create_engine

directory = '../HistoricalData'
config_file = "../Input/Config.xlsx"

from datetime import datetime
excel_date = 44575
dt = datetime.fromordinal(datetime(1900, 1, 1).toordinal() + excel_date - 2)
tt = dt.timetuple()
# print(dt)
# print(tt)


def extract_data_from_pptx(filename, eid):
    prs = Presentation(filename)
    output_filename = filename.split("\\")[1].split(".")[0]
    with pd.ExcelWriter(output_filename + '.xlsx', engine='openpyxl') as writer: 
        for i, slide in enumerate(prs.slides, 1):
            print(i)
            for shape in slide.shapes:
                # print(shape.has_chart)
                if shape.has_chart:
                    # print("yes")
                    # for p in shape.chart.plots:
                    #     # print(p)
                    #     print([c.label for c in p.categories])
                    # data = shape.chart.part.chart_workbook.xlsx_part.blob
                    # df = pd.read_excel(data)
                    # df.to_excel(output_filename+'.xlsx',str(i), index=False)

    
                    data = shape.chart.part.chart_workbook.xlsx_part.blob
                    df = pd.read_excel(data)
                    # df.rename(columns = {'Unnamed: 0':'TEST'}, inplace = True)
                    if i == 2:
                        cols = [
                            'RunDate',
                            'ContractMatch_Rate',
                            'ContractMatch_Rate_DistributorRemoved',
                            'Provider_Exception_Rate',
                            'Suppiler_Exception_Rate',
                            'Starting_Rate',
                            'Target',
                            ]
                        df.columns = cols
                    elif i == 3:
                        cols = [
                        'RunDate',
                        'ContractMatch_Rate',
                        'ContractMatch_Rate_DistributorRemoved',
                        'Starting_Rate',
                        'Target',
                        ]
                        df.columns = cols
                    elif i == 4:
                        cols = [
                        'RunDate',
                        'ContractMatch_Rate',
                        'ContractMatch_Rate_DistributorRemoved',
                        'Starting_Rate',
                        'Target',
                        ]
                        df.columns = cols
                    elif i == 5:
                        cols = [
                        'RunDate',
                        'ContractMatch_Rate',
                        'ContractMatch_Rate_DistributorRemoved',
                        'Starting_Rate'
                        ]
                        df.columns = cols
                    
                    # df.columns = cols
                    df['EID'] = eid
                    df['Slide'] = i
                    # print(df.columns)
                    df.to_excel(writer, sheet_name='Slide_' + str(i), index=False)
                        




                    # df.to_excel(output_filename+'.xlsx',"A", index=False)
                    # wb_temp = openpyxl.load_workbook(filename=BytesIO(data), data_only=True)
                    
                    # wb = openpyxl.Workbook()
                    # # ws = wb[(wb_temp.sheetnames[0])]
                    # ws = wb.create_sheet(str(i))
                    # wb.move_sheet()
                    # # ws.title = str(i)
                    # wb.save(output_filename+'.xlsx')

        
                    engine = create_engine("mysql://newuser:123456789@localhost/ppt_generation")
                    con = engine.connect()
                    # print(df.columns)
                    df.to_sql("historical_data",con,if_exists='append', index=False)
                    con.close()

                
        
config_df = pd.read_excel(config_file, "Config")



for filename in os.listdir(directory):
    f = os.path.join(directory, filename)
    if os.path.isfile(f):
        print(f)
        filtered_config_df = config_df[config_df['PPT_DisplayName']==filename]
        print(filtered_config_df['EID'].iloc[0])
        extract_data_from_pptx(f,filtered_config_df['EID'].iloc[0])


