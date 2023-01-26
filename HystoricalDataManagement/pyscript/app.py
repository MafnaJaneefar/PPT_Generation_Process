import os
import pandas as pd
import sys
from pptx import Presentation
from sqlalchemy import create_engine


#directory = "../HistoricalData"
#config_file = "../Input/Config.xlsx"
#output_dir = "../Output"

def extract_data_from_pptx(filename, eid):
    prs = Presentation(filename)
    output_filename = os.path.splitext(os.path.basename(filename))[0]

    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_chart:

                data = shape.chart.part.chart_workbook.xlsx_part.blob
                df = pd.read_excel(data)
                if i == 2:
                    cols = [
                        'RunDate',
                        'ContractExceptionRate',
                        'ContractExceptionRate_DistributorRemoved',
                        'ProviderException_Rate',
                        'SuppilerException_Rate',
                        'Starting_Rate',
                        'Target',
                        ]
                    df.columns = cols
                elif i == 3:
                    cols = [
                    'RunDate',
                    'ContractMatch_Rate',
                    'ContractMatch_Rate_DistributorRemoved',
                    'Starting_Rate',
                    'Target',
                    ]
                    df.columns = cols
                elif i == 4:
                    cols = [
                    'RunDate',
                    'ContractMatch_Rate',
                    'ContractMatch_Rate_DistributorRemoved',
                    'Starting_Rate',
                    'Target',
                    ]
                    df.columns = cols
                elif i == 5:
                    cols = [
                    'RunDate',
                    'ContractExceptionRate',
                    'ContractExceptionRate_DistributorRemoved',
                    'Starting_Rate'
                    ]
                    df.columns = cols

                df['EID'] = eid
                df['Slide'] = i
                df['Status'] = 'Active'

                engine = create_engine("mysql://newuser:123456789@localhost/ppt_generation")
                con = engine.connect()
                df.to_sql("historical_data",con,if_exists='append', index=False)
                con.close()


def historical_data_management(directory, config_file, output_dir):  
    config_df = pd.read_excel(config_file, "Config")
    status = {}
    for filename in os.listdir(directory):
        f = os.path.join(directory, filename)
        if os.path.isfile(f):
            print(f)
            filtered_config_df = config_df[config_df['PPT_DisplayName']==filename]
            # print(filtered_config_df['EID'].iloc[0])
            try:
                extract_data_from_pptx(f, filtered_config_df['EID'].iloc[0])
                data_inserted = True
            except Exception as e:
                data_inserted = False
                print(e)
            
            status[os.path.basename(filename)] = data_inserted

    excel_output = output_dir + "/process_status.xlsx"
    status_df = pd.DataFrame(list(status.items()), columns=["File Name", "Status"])
    status_df.to_excel(excel_output, index=False)



historical_data_management(sys.argv[1],sys.argv[2],sys.argv[3])

#historical_data_management(directory, config_file, output_dir)
